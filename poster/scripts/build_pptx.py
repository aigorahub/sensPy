"""Assemble the sensPy Sensometrics 2026 poster PowerPoint."""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
POSTER = ROOT / "poster"
CHARTS = POSTER / "charts"
CHART_DATA = POSTER / "chart_data"
ASSETS = POSTER / "assets"
OUT = POSTER / "print_artifacts"
OUT.mkdir(parents=True, exist_ok=True)

PPTX = OUT / "senspy-sensometrics-2026-poster.pptx"

SLIDE_W = 27.83
SLIDE_H = 39.37
M = 0.85
HEADER_H = 4.35
FOOTER_H = 1.45
GAP = 0.55
COL_W = (SLIDE_W - 2 * M - GAP) / 2
LEFT_X = M
RIGHT_X = M + COL_W + GAP

DISPLAY = "Georgia"
BODY = "Arial"
MONO = "Courier New"

CANVAS = "f4f0e6"
PANEL = "fbfaf6"
INK = "17291f"
CORAL = "c7563f"
GREEN = "4c8a61"
SAGE = "9aa79b"
MIST = "e8e3d8"
DARK = "10231a"
WHITE = "ffffff"


def rgb(hex_value: str) -> RGBColor:
    h = hex_value.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_rect(slide, x, y, w, h, fill, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    rect = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = rgb(fill)
    if line:
        rect.line.color.rgb = rgb(line)
        rect.line.width = Pt(1)
    else:
        rect.line.fill.background()
    return rect


def add_text(
    slide,
    text: str,
    x,
    y,
    w,
    h,
    *,
    size=18,
    font=BODY,
    color=INK,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.05,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.font.name = font
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(4)
    return box


def add_heading(slide, title: str, x, y, w, *, size=29):
    add_rect(slide, x, y, w, 0.055, CORAL)
    add_text(
        slide,
        title,
        x,
        y + 0.12,
        w,
        0.55,
        size=size,
        font=DISPLAY,
        color=DARK,
        align=PP_ALIGN.CENTER,
        margin=0,
    )
    return y + 0.78


def add_image_fit(slide, path: Path, x, y, w, h):
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    pic_h = pic.height / 914400.0
    pic_w = pic.width / 914400.0
    if pic_h > h:
        scale = h / pic_h
        pic.width = int(pic.width * scale)
        pic.height = int(pic.height * scale)
        pic_w = pic.width / 914400.0
        pic_h = pic.height / 914400.0
    pic.left = Inches(x + (w - pic_w) / 2)
    pic.top = Inches(y + (h - pic_h) / 2)
    return pic


def add_metric(slide, x, y, w, h, value, label, sub=None, *, color=CORAL):
    add_rect(slide, x, y, w, h, MIST, line="#d2c9bb", radius=True)
    add_text(
        slide,
        value,
        x + 0.08,
        y + 0.14,
        w - 0.16,
        0.55,
        size=34,
        font=BODY,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0,
    )
    add_text(
        slide,
        label,
        x + 0.12,
        y + 0.83,
        w - 0.24,
        0.5,
        size=16,
        font=BODY,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0,
    )
    if sub:
        add_text(
            slide,
            sub,
            x + 0.14,
            y + 1.28,
            w - 0.28,
            0.48,
            size=13.5,
            font=BODY,
            color="#4d5a52",
            align=PP_ALIGN.CENTER,
            margin=0,
        )


def add_compact_metric(slide, x, y, w, h, value, label, sub=None, *, color=CORAL):
    add_rect(slide, x, y, w, h, MIST, line="#d2c9bb", radius=True)
    add_text(
        slide,
        value,
        x + 0.06,
        y + 0.08,
        w - 0.12,
        0.52,
        size=30,
        font=BODY,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0,
    )
    add_text(
        slide,
        label,
        x + 0.1,
        y + 0.73,
        w - 0.2,
        0.4,
        size=14.2,
        font=BODY,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0,
    )
    if sub:
        add_text(
            slide,
            sub,
            x + 0.12,
            y + 1.17,
            w - 0.24,
            h - 1.17,
            size=12.5,
            font=BODY,
            color="#4d5a52",
            align=PP_ALIGN.CENTER,
            margin=0,
        )


def add_callout(slide, x, y, w, h, title, body):
    add_rect(slide, x, y, w, h, MIST, line="#d2c9bb", radius=True)
    add_text(slide, title, x + 0.22, y + 0.18, w - 0.44, 0.35, size=17, color=CORAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, body, x + 0.32, y + 0.62, w - 0.64, h - 0.72, size=16.3, color=INK, align=PP_ALIGN.CENTER)


def add_takeaway(slide, x, y, w, h, title, body, *, color=CORAL):
    add_rect(slide, x, y, w, h, MIST, line="#d2c9bb", radius=True)
    add_text(slide, title, x + 0.22, y + 0.18, w - 0.44, 0.4, size=18, color=color, bold=True, align=PP_ALIGN.CENTER, margin=0)
    add_text(slide, body, x + 0.34, y + 0.76, w - 0.68, h - 0.9, size=15.3, color=INK, align=PP_ALIGN.CENTER, margin=0)


def build_header(slide):
    add_rect(slide, 0, 0, SLIDE_W, HEADER_H, MIST)
    add_rect(slide, 0, HEADER_H - 0.07, SLIDE_W, 0.07, CORAL)
    add_text(
        slide,
        "sensPy",
        M,
        0.5,
        SLIDE_W - 2 * M,
        1.25,
        size=88,
        font=DISPLAY,
        color=DARK,
        align=PP_ALIGN.CENTER,
        margin=0,
    )
    add_text(
        slide,
        "Bringing the gold standard of sensR to the Python sensory ecosystem",
        M,
        1.78,
        SLIDE_W - 2 * M,
        0.82,
        size=32,
        font=BODY,
        color=CORAL,
        align=PP_ALIGN.CENTER,
        margin=0,
    )
    add_text(
        slide,
        "John M. Ennis | Aigora | Sensometrics 2026",
        M,
        3.0,
        SLIDE_W - 2 * M,
        0.5,
        size=22,
        font=BODY,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0,
    )
    add_text(slide, "Python-native Thurstonian sensory discrimination", 20.9, 0.62, 5.4, 0.5, size=15.5, color="#526057", align=PP_ALIGN.RIGHT)


def build_footer(slide, summary):
    y = SLIDE_H - FOOTER_H
    add_rect(slide, 0, y, SLIDE_W, FOOTER_H, DARK)
    add_text(slide, "sensPy", M, y + 0.42, 2.0, 0.45, size=20, font=DISPLAY, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(
        slide,
        "github.com/aigorahub/sensPy | package v" + summary["version"],
        8.1,
        y + 0.44,
        11.5,
        0.45,
        size=17,
        font=BODY,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    qr = ASSETS / "qr-senspy-github.png"
    if qr.exists():
        slide.shapes.add_picture(str(qr), Inches(SLIDE_W - M - 0.92), Inches(y + 0.24), width=Inches(0.92))


def build_problem(slide):
    y = add_heading(slide, "The problem", LEFT_X, 4.7, COL_W)
    add_text(
        slide,
        "sensR is the trusted R reference for sensory discrimination. Python teams need the same signal-detection discipline without leaving their SciPy, dataclass, and Plotly workflows.",
        LEFT_X + 0.05,
        y + 0.05,
        COL_W - 0.1,
        1.15,
        size=21,
        color=INK,
        align=PP_ALIGN.CENTER,
    )
    add_callout(
        slide,
        LEFT_X + 0.2,
        y + 1.5,
        (COL_W - 0.6) / 2,
        1.45,
        "Gold standard",
        "Preserve the sensR numerical contract.",
    )
    add_callout(
        slide,
        LEFT_X + 0.4 + (COL_W - 0.6) / 2,
        y + 1.5,
        (COL_W - 0.6) / 2,
        1.45,
        "Python-native",
        "Typed objects, plots, and planning tools.",
    )
    add_text(
        slide,
        "Goal: make sensory discrimination methods feel native in Python without loosening the validation discipline.",
        LEFT_X + 0.15,
        y + 3.28,
        COL_W - 0.3,
        0.8,
        size=19,
        color=CORAL,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def build_metrics(slide, summary):
    y = add_heading(slide, "What sensPy ports", RIGHT_X, 4.7, COL_W)
    card_w = (COL_W - 0.55) / 3
    xs = [RIGHT_X + 0.05, RIGHT_X + 0.275 + card_w, RIGHT_X + 0.5 + 2 * card_w]
    add_compact_metric(slide, xs[0], y + 0.1, card_w, 1.45, "13", "protocol variants", "8 single + 5 double")
    add_compact_metric(slide, xs[1], y + 0.1, card_w, 1.45, "740+", "automated tests", f"{summary['collected_pytest_items']} collected cases", color=GREEN)
    add_compact_metric(slide, xs[2], y + 0.1, card_w, 1.45, summary["sensr_version"], "sensR fixture", "golden parity data")
    add_compact_metric(slide, xs[0], y + 1.78, card_w, 1.45, str(summary["dataclass_count"]), "dataclasses", "typed Python results", color=GREEN)
    add_compact_metric(slide, xs[1], y + 1.78, card_w, 1.45, str(summary["api_export_count"]), "public exports", "single import surface", color=GREEN)
    add_compact_metric(slide, xs[2], y + 1.78, card_w, 1.45, "Plotly", "interactive figures", "psychometric + ROC")
    add_text(
        slide,
        "Numerical parity targets: 3-6 decimal places across sensR-backed fixtures and boundary cases.",
        RIGHT_X + 0.18,
        y + 3.56,
        COL_W - 0.36,
        0.55,
        size=17.5,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def build_charts(slide):
    y = add_heading(slide, "Protocol coverage", LEFT_X, 9.65, COL_W)
    add_image_fit(slide, CHARTS / "protocol_coverage.png", LEFT_X + 0.1, y + 0.02, COL_W - 0.2, 4.75)
    add_text(slide, "Common d-prime scale, protocol-specific guessing.", LEFT_X + 0.25, y + 4.9, COL_W - 0.5, 0.55, size=16, color="#4d5a52", align=PP_ALIGN.CENTER)

    y2 = add_heading(slide, "Psychometric functions", RIGHT_X, 9.65, COL_W)
    add_image_fit(slide, CHARTS / "psychometric_curves.png", RIGHT_X + 0.0, y2 + 0.02, COL_W, 4.78)
    add_text(slide, "One API maps Pc, Pd, and d-prime.", RIGHT_X + 0.25, y2 + 4.92, COL_W - 0.5, 0.5, size=16, color="#4d5a52", align=PP_ALIGN.CENTER)

    y3 = add_heading(slide, "Validation surface", LEFT_X, 16.25, COL_W)
    add_image_fit(slide, CHARTS / "test_inventory.png", LEFT_X + 0.0, y3 + 0.02, COL_W, 4.95)
    add_text(slide, "Golden sensR fixtures run beside unit, model, power, and plotting tests.", LEFT_X + 0.25, y3 + 5.08, COL_W - 0.5, 0.5, size=16, color="#4d5a52", align=PP_ALIGN.CENTER)

    y4 = add_heading(slide, "SciPy-native architecture", RIGHT_X, 16.25, COL_W)
    add_image_fit(slide, CHARTS / "architecture_pipeline.png", RIGHT_X + 0.0, y4 + 0.06, COL_W, 2.45)
    add_image_fit(slide, CHARTS / "roc_bridge.png", RIGHT_X + 0.35, y4 + 2.7, COL_W - 0.7, 3.1)


def build_api_and_models(slide):
    y = add_heading(slide, "Python ergonomics", LEFT_X, 22.95, COL_W)
    add_rect(slide, LEFT_X + 0.2, y + 0.05, COL_W - 0.4, 3.2, MIST, line="#d2c9bb", radius=True)
    code = (
        "from senspy import discrim, dprime_power\n"
        "\n"
        "result = discrim(correct=80, total=100,\n"
        "                 method=\"triangle\")\n"
        "result.d_prime\n"
        "result.confint()\n"
        "\n"
        "dprime_power(1.5, 100,\n"
        "             method=\"triangle\")"
    )
    add_text(slide, code, LEFT_X + 0.55, y + 0.28, COL_W - 1.1, 2.7, size=15.2, font=MONO, color=INK)
    add_text(
        slide,
        "Typed results keep estimates, intervals, p-values, and summaries together.",
        LEFT_X + 0.35,
        y + 3.42,
        COL_W - 0.7,
        0.8,
        size=18,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    y2 = add_heading(slide, "Advanced models", RIGHT_X, 22.95, COL_W)
    w = (COL_W - 0.75) / 3
    add_metric(slide, RIGHT_X + 0.1, y2 + 0.2, w, 1.7, "BB", "Beta-Binomial", "overdispersed panel data", color=GREEN)
    add_metric(slide, RIGHT_X + 0.35 + w, y2 + 0.2, w, 1.7, "SD", "Same-Different", "Thurstonian criterion model")
    add_metric(slide, RIGHT_X + 0.6 + 2 * w, y2 + 0.2, w, 1.7, "ROC", "SDT analysis", "AUC + rating data", color=GREEN)
    add_text(
        slide,
        "Also included: 2-AC, DOD, A-not-A, d-prime tests, posthoc, simulation, power, and sample-size planning.",
        RIGHT_X + 0.25,
        y2 + 2.25,
        COL_W - 0.5,
        0.8,
        size=17.5,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def build_conclusion(slide):
    y = 28.25
    add_heading(slide, "Conclusion", M, y, SLIDE_W - 2 * M)
    add_text(
        slide,
        "sensPy brings the sensR contract into Python: the same sensory discrimination ideas, validated against the R reference, expressed as SciPy-native functions, typed dataclasses, and interactive Plotly figures.",
        M + 1.1,
        y + 0.86,
        SLIDE_W - 2 * M - 2.2,
        1.15,
        size=24,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "Selected references: Brockhoff & Christensen, sensR package; Macmillan & Creelman, Detection Theory; SciPy; Plotly. Repository includes golden sensR v1.5.3 fixtures.",
        M + 1.6,
        y + 2.12,
        SLIDE_W - 2 * M - 3.2,
        0.7,
        size=15.5,
        color="#4d5a52",
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "What changes for sensometricians",
        M,
        y + 2.95,
        SLIDE_W - 2 * M,
        0.42,
        size=21,
        font=DISPLAY,
        color=DARK,
        align=PP_ALIGN.CENTER,
        margin=0,
    )
    take_w = (SLIDE_W - 2 * M - 0.8) / 3
    add_takeaway(
        slide,
        M,
        y + 3.42,
        take_w,
        2.05,
        "Validated migration",
        "Keep sensR-backed decisions in Python notebooks and pipelines.",
        color=GREEN,
    )
    add_takeaway(
        slide,
        M + take_w + 0.4,
        y + 3.42,
        take_w,
        2.05,
        "One reporting scale",
        "Pc, Pd, d-prime, intervals, and power share one API.",
    )
    add_takeaway(
        slide,
        M + 2 * (take_w + 0.4),
        y + 3.42,
        take_w,
        2.05,
        "Modern outputs",
        "Typed results and Plotly figures are easier to audit and reuse.",
        color=GREEN,
    )
    strip_y = 33.95
    add_text(
        slide,
        "Open Python sensory workflow",
        M,
        strip_y,
        SLIDE_W - 2 * M,
        0.42,
        size=20,
        font=DISPLAY,
        color=DARK,
        align=PP_ALIGN.CENTER,
        margin=0,
    )
    card_w = (SLIDE_W - 2 * M - 1.2) / 4
    cards = [
        ("1", "Estimate", "d-prime, Pc, Pd, intervals", "discrim(), samediff(), dod()", "typed estimates + confidence intervals"),
        ("2", "Plan", "power and sample size", "discrim_power(), dprime_sample_size()", "design studies before collecting data"),
        ("3", "Model", "overdispersion and criteria", "betabin(), samediff(), twoac()", "panel variability + criterion effects"),
        ("4", "Visualize", "Plotly psychometric + ROC", "plot_psychometric(), plot_roc()", "interactive curves for review"),
    ]
    for idx, (num, title, body, api, output) in enumerate(cards):
        x = M + idx * (card_w + 0.4)
        add_rect(slide, x, strip_y + 0.58, card_w, 3.05, MIST, line="#d2c9bb", radius=True)
        add_text(slide, num + ". " + title, x + 0.15, strip_y + 0.82, card_w - 0.3, 0.42, size=20, color=CORAL, bold=True, align=PP_ALIGN.CENTER, margin=0)
        add_text(slide, body, x + 0.42, strip_y + 1.42, card_w - 0.84, 0.6, size=16, color=INK, align=PP_ALIGN.CENTER, margin=0)
        add_text(slide, api, x + 0.42, strip_y + 2.14, card_w - 0.84, 0.44, size=14.2, font=MONO, color=DARK, align=PP_ALIGN.CENTER, margin=0)
        add_text(slide, output, x + 0.42, strip_y + 2.68, card_w - 0.84, 0.48, size=13.8, color="#4d5a52", align=PP_ALIGN.CENTER, margin=0)


def main() -> None:
    summary = json.loads((CHART_DATA / "summary.json").read_text())
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, CANVAS)
    build_header(slide)
    build_problem(slide)
    build_metrics(slide, summary)
    build_charts(slide)
    build_api_and_models(slide)
    build_conclusion(slide)
    build_footer(slide, summary)

    prs.save(PPTX)
    print(f"[pptx] wrote {PPTX.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
